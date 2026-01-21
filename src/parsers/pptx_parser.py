from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from .base_parser import BaseParser


class PPTXParser(BaseParser):
    """Parser for PPTX documents"""

    def parse(self) -> str:
        """Extract content from PowerPoint and convert to markdown"""
        markdown_lines = []
        markdown_lines.append(f"# {self.file_path.stem}\n")
        markdown_lines.append(f"*Source: {self.file_path.name}*\n")
        markdown_lines.append("---\n")

        try:
            prs = Presentation(self.file_path)

            for slide_num, slide in enumerate(prs.slides, 1):
                # Add slide header
                markdown_lines.append(f"\n## Slide {slide_num}\n")

                # Process slide title
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                    if title:
                        markdown_lines.append(f"### {title}\n")

                # Process all shapes in the slide
                for shape in slide.shapes:
                    # Text boxes and placeholders
                    if hasattr(shape, "text") and shape.text.strip():
                        # Skip if it's the title (already processed)
                        if shape != slide.shapes.title:
                            text = shape.text.strip()
                            markdown_lines.append(f"{text}\n")

                    # Tables
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        markdown_lines.append(self._table_to_markdown(shape.table))
                        markdown_lines.append("\n")

                # Process notes if any
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        markdown_lines.append(f"\n**Notes:**\n{notes_text}\n")

                markdown_lines.append("\n---\n")

        except Exception as e:
            markdown_lines.append(f"\n**Error parsing PPTX:** {str(e)}\n")

        return "\n".join(markdown_lines)

    def _table_to_markdown(self, table) -> str:
        """Convert PowerPoint table to markdown format"""
        if not table.rows:
            return ""

        markdown_table = []

        # Header row
        header_row = table.rows[0]
        header = [cell.text.strip() for cell in header_row.cells]
        markdown_table.append("| " + " | ".join(header) + " |")
        markdown_table.append("| " + " | ".join(["---"] * len(header)) + " |")

        # Data rows
        for row in table.rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            markdown_table.append("| " + " | ".join(cells) + " |")

        return "\n".join(markdown_table)
